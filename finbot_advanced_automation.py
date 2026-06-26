"""
FinBot AI Advanced Automation Dashboard
Python + NumPy + pandas + matplotlib + seaborn + Streamlit
Excel + Word + PowerPoint automation

Run:
pip install streamlit pandas numpy matplotlib seaborn openpyxl python-docx python-pptx xlsxwriter
streamlit run finbot_advanced_automation.py
"""

import io
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import streamlit as st

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PPTInches


st.set_page_config(page_title="FinBot AI Advanced Automation", page_icon="🤖", layout="wide")

st.title("🤖 FinBot AI Advanced Automation Dashboard")
st.write("Excel analysis + Word report + PowerPoint deck + Python charts + analyst insights")


def sample_data():
    return pd.DataFrame({
        "Year": [2021, 2022, 2023, 2024, 2025],
        "Revenue": [1200000, 1450000, 1680000, 2100000, 2520000],
        "EBITDA": [240000, 326000, 420000, 588000, 756000],
        "NetProfit": [120000, 178000, 240000, 345000, 478000],
        "Assets": [2500000, 2800000, 3300000, 3900000, 4700000],
        "Liabilities": [900000, 980000, 1200000, 1450000, 1600000],
    })


def load_data(uploaded_file):
    if uploaded_file is None:
        return sample_data()

    name = uploaded_file.name.lower()
    if name.endswith(".csv"):
        return pd.read_csv(uploaded_file)
    if name.endswith((".xlsx", ".xls")):
        return pd.read_excel(uploaded_file)

    st.error("Upload CSV or Excel file only.")
    return sample_data()


def clean_data(df):
    df = df.copy()
    df.columns = [str(c).strip() for c in df.columns]

    mapping = {
        "year": "Year",
        "date": "Year",
        "sales": "Revenue",
        "revenue": "Revenue",
        "ebitda": "EBITDA",
        "profit": "NetProfit",
        "net profit": "NetProfit",
        "netprofit": "NetProfit",
        "assets": "Assets",
        "liabilities": "Liabilities",
    }

    df = df.rename(columns={c: mapping.get(c.lower(), c) for c in df.columns})

    for col in ["Year", "Revenue", "EBITDA", "NetProfit"]:
        if col not in df.columns:
            df[col] = 0

    for col in ["Revenue", "EBITDA", "NetProfit", "Assets", "Liabilities"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    df = df.sort_values("Year").reset_index(drop=True)

    df["Revenue_Growth_%"] = df["Revenue"].pct_change().replace([np.inf, -np.inf], 0).fillna(0) * 100
    df["EBITDA_Margin_%"] = np.where(df["Revenue"] != 0, df["EBITDA"] / df["Revenue"] * 100, 0)
    df["Net_Profit_Margin_%"] = np.where(df["Revenue"] != 0, df["NetProfit"] / df["Revenue"] * 100, 0)

    if "Assets" in df.columns and "Liabilities" in df.columns:
        df["Equity"] = df["Assets"] - df["Liabilities"]
        df["Debt_to_Equity"] = np.where(df["Equity"] != 0, df["Liabilities"] / df["Equity"], 0)

    return df


def calculate_kpis(df):
    total_revenue = df["Revenue"].sum()
    total_ebitda = df["EBITDA"].sum()
    total_profit = df["NetProfit"].sum()

    first_revenue = df["Revenue"].iloc[0] if len(df) else 0
    last_revenue = df["Revenue"].iloc[-1] if len(df) else 0

    growth = ((last_revenue - first_revenue) / first_revenue * 100) if first_revenue else 0
    ebitda_margin = (total_ebitda / total_revenue * 100) if total_revenue else 0
    profit_margin = (total_profit / total_revenue * 100) if total_revenue else 0

    periods = max(len(df) - 1, 1)
    cagr = ((last_revenue / first_revenue) ** (1 / periods) - 1) * 100 if first_revenue else 0

    if growth > 50 and ebitda_margin > 25:
        rating = "Strong Buy"
    elif growth > 20 and ebitda_margin > 15:
        rating = "Positive"
    elif growth > 0:
        rating = "Watchlist"
    else:
        rating = "Risk"

    return {
        "Total Revenue": total_revenue,
        "Total EBITDA": total_ebitda,
        "Total Net Profit": total_profit,
        "Revenue Growth %": growth,
        "Revenue CAGR %": cagr,
        "EBITDA Margin %": ebitda_margin,
        "Net Profit Margin %": profit_margin,
        "Rating": rating,
    }


def analyst_summary(kpis):
    return f"""
FinBot AI Analyst Report

Business Rating: {kpis['Rating']}

Key KPIs:
Total Revenue: ₹{kpis['Total Revenue']:,.0f}
Total EBITDA: ₹{kpis['Total EBITDA']:,.0f}
Total Net Profit: ₹{kpis['Total Net Profit']:,.0f}
Revenue Growth: {kpis['Revenue Growth %']:.2f}%
Revenue CAGR: {kpis['Revenue CAGR %']:.2f}%
EBITDA Margin: {kpis['EBITDA Margin %']:.2f}%
Net Profit Margin: {kpis['Net Profit Margin %']:.2f}%

Analyst View:
Revenue performance is {"strong" if kpis["Revenue Growth %"] > 30 else "moderate"}.
Operating margin is {"healthy" if kpis["EBITDA Margin %"] > 20 else "needs improvement"}.
Profitability is {"good" if kpis["Net Profit Margin %"] > 10 else "weak"}.

Investment Banking Use:
This output can support DCF valuation, comparable company analysis,
pitch deck preparation, management presentation, and equity research reports.
"""


def create_revenue_chart(df):
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.plot(df["Year"], df["Revenue"], marker="o", label="Revenue")
    ax.plot(df["Year"], df["EBITDA"], marker="o", label="EBITDA")
    ax.plot(df["Year"], df["NetProfit"], marker="o", label="Net Profit")
    ax.set_title("Revenue, EBITDA and Net Profit Trend")
    ax.set_xlabel("Year")
    ax.set_ylabel("Amount")
    ax.legend()
    ax.grid(True)
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", bbox_inches="tight")
    buffer.seek(0)
    plt.close(fig)
    return buffer


def create_heatmap(df):
    numeric_df = df.select_dtypes(include=[np.number])
    fig, ax = plt.subplots(figsize=(8, 5))
    sns.heatmap(numeric_df.corr(), annot=True, cmap="Blues", ax=ax)
    ax.set_title("Financial Correlation Heatmap")
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", bbox_inches="tight")
    buffer.seek(0)
    plt.close(fig)
    return buffer


def export_excel(df, kpis):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Clean_Data", index=False)
        pd.DataFrame([kpis]).to_excel(writer, sheet_name="KPI_Summary", index=False)

        workbook = writer.book
        worksheet = writer.sheets["Clean_Data"]
        money_fmt = workbook.add_format({"num_format": "₹#,##0"})
        pct_fmt = workbook.add_format({"num_format": "0.00%"})

        worksheet.set_column("A:A", 12)
        worksheet.set_column("B:D", 18, money_fmt)
        worksheet.set_column("E:G", 18)

    output.seek(0)
    return output


def export_word_report(df, kpis, summary):
    doc = Document()
    doc.add_heading("FinBot AI Analyst Report", 0)

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(summary)

    doc.add_heading("KPI Summary", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"

    for key, value in kpis.items():
        cells = table.add_row().cells
        cells[0].text = str(key)
        if isinstance(value, (int, float)):
            cells[1].text = f"{value:,.2f}"
        else:
            cells[1].text = str(value)

    chart_buffer = create_revenue_chart(df)
    doc.add_heading("Financial Trend Chart", level=1)
    doc.add_picture(chart_buffer, width=Inches(6))

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def export_powerpoint(df, kpis, summary):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FinBot AI Financial Analysis"
    slide.placeholders[1].text = "Excel + Word + PowerPoint Automation Dashboard"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPI Summary"
    body = slide.placeholders[1]
    body.text = (
        f"Total Revenue: ₹{kpis['Total Revenue']:,.0f}\n"
        f"Revenue Growth: {kpis['Revenue Growth %']:.2f}%\n"
        f"Revenue CAGR: {kpis['Revenue CAGR %']:.2f}%\n"
        f"EBITDA Margin: {kpis['EBITDA Margin %']:.2f}%\n"
        f"Net Profit Margin: {kpis['Net Profit Margin %']:.2f}%\n"
        f"Rating: {kpis['Rating']}"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Financial Trend Chart"
    chart_buffer = create_revenue_chart(df)
    slide.shapes.add_picture(chart_buffer, PPTInches(1), PPTInches(1.4), width=PPTInches(8))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Analyst Recommendation"
    slide.placeholders[1].text = summary[:1200]

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


uploaded_file = st.sidebar.file_uploader("Upload CSV / Excel", type=["csv", "xlsx", "xls"])

df_raw = load_data(uploaded_file)
df = clean_data(df_raw)
kpis = calculate_kpis(df)
summary = analyst_summary(kpis)

c1, c2, c3, c4 = st.columns(4)
c1.metric("Total Revenue", f"₹{kpis['Total Revenue']:,.0f}")
c2.metric("Revenue Growth", f"{kpis['Revenue Growth %']:.2f}%")
c3.metric("EBITDA Margin", f"{kpis['EBITDA Margin %']:.2f}%")
c4.metric("Rating", kpis["Rating"])

st.subheader("Clean Financial Data")
st.dataframe(df, use_container_width=True)

st.subheader("Charts")
left, right = st.columns(2)

with left:
    st.write("Financial Trend")
    chart_buffer = create_revenue_chart(df)
    st.image(chart_buffer)

with right:
    st.write("Correlation Heatmap")
    heatmap_buffer = create_heatmap(df)
    st.image(heatmap_buffer)

st.subheader("FinBot Analyst Summary")
st.text_area("Generated Report", value=summary, height=300)

st.subheader("Download Automation Outputs")

excel_file = export_excel(df, kpis)
word_file = export_word_report(df, kpis, summary)
ppt_file = export_powerpoint(df, kpis, summary)

d1, d2, d3 = st.columns(3)

with d1:
    st.download_button(
        "Download Excel Automation",
        data=excel_file,
        file_name="FinBot_Excel_Automation.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

with d2:
    st.download_button(
        "Download Word Report",
        data=word_file,
        file_name="FinBot_Word_Analyst_Report.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

with d3:
    st.download_button(
        "Download PowerPoint Deck",
        data=ppt_file,
        file_name="FinBot_PowerPoint_Deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

st.subheader("FinBot Prompt Automation")
prompt = st.text_input("Ask FinBot", "Prepare investment banking analyst summary")

if st.button("Run FinBot"):
    st.success("Automation completed successfully.")
    st.write(summary)
